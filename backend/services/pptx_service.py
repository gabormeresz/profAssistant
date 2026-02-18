"""
PowerPoint (PPTX) generation service.

Converts a Presentation schema into a styled .pptx file using python-pptx.
Clean, neutral design — easy for teachers to restyle with their own branding.
"""

from __future__ import annotations

import io
import logging
from typing import TYPE_CHECKING, Optional

from pptx import Presentation as make_pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PptxPresentationType

from schemas.presentation import Presentation as PresentationModel, Slide as SlideModel

logger = logging.getLogger(__name__)

# ── Neutral colour palette ──────────────────────────────────────────────────
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)  # very light bg tint
_GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)  # separator line
_GRAY_500 = RGBColor(0x6B, 0x72, 0x80)  # secondary text
_GRAY_800 = RGBColor(0x1F, 0x2A, 0x37)  # primary text
_GRAY_900 = RGBColor(0x11, 0x18, 0x27)  # title text


# ── Helpers ─────────────────────────────────────────────────────────────────


def _add_textbox(slide_obj, left, top, width, height):
    """Shortcut that returns (shape, text_frame)."""
    shape = slide_obj.shapes.add_textbox(left, top, width, height)
    return shape, shape.text_frame


def _set_paragraph(
    text_frame,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = _GRAY_800,
    alignment=PP_ALIGN.LEFT,
    space_before: Optional[int] = None,
    space_after: Optional[int] = None,
    font_name: str = "Calibri",
):
    """Write *text* into the first paragraph of *text_frame*."""
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.text = text
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.font.name = font_name
    para.alignment = alignment
    if space_before is not None:
        para.space_before = Pt(space_before)
    if space_after is not None:
        para.space_after = Pt(space_after)


def _add_paragraph(
    text_frame,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = _GRAY_800,
    alignment=PP_ALIGN.LEFT,
    space_before: Optional[int] = None,
    space_after: Optional[int] = None,
    font_name: str = "Calibri",
):
    """Append a new paragraph to *text_frame*."""
    para = text_frame.add_paragraph()
    para.text = text
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.font.name = font_name
    para.alignment = alignment
    if space_before is not None:
        para.space_before = Pt(space_before)
    if space_after is not None:
        para.space_after = Pt(space_after)
    return para


def _add_slide_number_footer(slide_obj, slide_number: int, total: int) -> None:
    """Small page-number label at the bottom-right corner."""
    _, tf = _add_textbox(
        slide_obj, Inches(8.5), Inches(5.15), Inches(1.2), Inches(0.35)
    )
    _set_paragraph(
        tf,
        f"{slide_number} / {total}",
        font_size=9,
        color=_GRAY_500,
        alignment=PP_ALIGN.RIGHT,
    )


# ── Slide builders ──────────────────────────────────────────────────────────


def _build_title_slide(
    prs: PptxPresentationType, presentation: PresentationModel
) -> None:
    """Clean white title slide with centred text."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide_obj = prs.slides.add_slide(slide_layout)

    # Course title — small subtitle above the main title
    _, tf = _add_textbox(slide_obj, Inches(1), Inches(1.6), Inches(8), Inches(0.6))
    _set_paragraph(
        tf,
        presentation.course_title,
        font_size=14,
        color=_GRAY_500,
        alignment=PP_ALIGN.CENTER,
    )

    # Lesson title — large, centred
    _, tf = _add_textbox(slide_obj, Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.6))
    _set_paragraph(
        tf,
        presentation.lesson_title,
        font_size=32,
        bold=True,
        color=_GRAY_900,
        alignment=PP_ALIGN.CENTER,
    )

    # Thin horizontal rule
    rule = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.85), Inches(2.0), Pt(1.5)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = _GRAY_300
    rule.line.fill.background()

    # Class & slide count
    _, tf = _add_textbox(slide_obj, Inches(3.0), Inches(4.1), Inches(4), Inches(0.5))
    _set_paragraph(
        tf,
        f"Class {presentation.class_number}  ·  {len(presentation.slides)} slides",
        font_size=12,
        color=_GRAY_500,
        alignment=PP_ALIGN.CENTER,
    )


def _build_content_slide(
    prs: PptxPresentationType,
    slide: SlideModel,
    total_slides: int,
) -> None:
    """Clean content slide — title, separator, bullets, notes."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide_obj = prs.slides.add_slide(slide_layout)

    # ── Title ──
    _, tf = _add_textbox(slide_obj, Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.6))
    _set_paragraph(tf, slide.title, font_size=24, bold=True, color=_GRAY_900)

    # Thin gray rule below title
    rule = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.95), Inches(8.8), Pt(1)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = _GRAY_300
    rule.line.fill.background()

    # ── Bullet points ──
    _, tf = _add_textbox(slide_obj, Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.4))
    tf.word_wrap = True

    for i, point in enumerate(slide.bullet_points):
        if i == 0:
            _set_paragraph(
                tf, f"•  {point}", font_size=17, color=_GRAY_800, space_after=6
            )
        else:
            _add_paragraph(
                tf,
                f"•  {point}",
                font_size=17,
                color=_GRAY_800,
                space_before=4,
                space_after=6,
            )

    # ── Visual suggestion (small italic note at bottom) ──
    if slide.visual_suggestion:
        _, tf = _add_textbox(
            slide_obj, Inches(0.6), Inches(4.7), Inches(8.0), Inches(0.5)
        )
        para = tf.paragraphs[0]
        para.text = f"Visual suggestion: {slide.visual_suggestion}"
        para.font.size = Pt(10)
        para.font.italic = True
        para.font.color.rgb = _GRAY_500
        para.font.name = "Calibri"

    # ── Slide number footer ──
    _add_slide_number_footer(slide_obj, slide.slide_number, total_slides)

    # ── Speaker notes ──
    if slide.speaker_notes:
        notes_slide = slide_obj.notes_slide
        notes_tf = notes_slide.notes_text_frame
        if notes_tf is not None:
            notes_tf.text = slide.speaker_notes


# ── Public API ──────────────────────────────────────────────────────────────


def generate_pptx(presentation: PresentationModel) -> bytes:
    """
    Build a .pptx file from a *PresentationModel* schema and return the raw bytes.

    Parameters
    ----------
    presentation : PresentationModel
        Validated Pydantic model with course_title, lesson_title,
        class_number, and a list of slides.

    Returns
    -------
    bytes
        The binary content of the generated .pptx file.
    """
    prs = make_pptx()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    _build_title_slide(prs, presentation)

    total = len(presentation.slides)
    for slide in presentation.slides:
        _build_content_slide(prs, slide, total)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    logger.info(
        "Generated PPTX: %s – %s (class %d, %d slides)",
        presentation.course_title,
        presentation.lesson_title,
        presentation.class_number,
        len(presentation.slides),
    )
    return buf.getvalue()
